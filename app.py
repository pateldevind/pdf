from flask import Flask, render_template, send_from_directory, abort, request, jsonify, send_file, session
import os
import logging
import datetime
from werkzeug.utils import secure_filename
from pdf2docx import Converter
import tempfile
import fitz  # PyMuPDF
from PIL import Image
import io
import zipfile
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4, LETTER, LEGAL
from docx import Document
from reportlab.lib.units import inch
from pptx import Presentation
from openpyxl import load_workbook
from reportlab.lib import colors
import requests
import pdfkit
import PyPDF2
import roman
import jwt
from functools import wraps
import sqlite3
from werkzeug.security import generate_password_hash, check_password_hash
from pdf2image import convert_from_path

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Get the absolute path of the current directory
BASE_DIR = os.path.abspath(os.path.dirname(__file__))

app = Flask(__name__, static_url_path='', static_folder='.')
app.secret_key = os.urandom(24)  # Generate a secure secret key
JWT_SECRET_KEY = os.urandom(24)  # Generate a secure JWT secret key

# Configure upload folder
UPLOAD_FOLDER = os.path.join(BASE_DIR, 'uploads')
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

# Define page sizes
PAGE_SIZES = {
    'a4': A4,
    'letter': LETTER,
    'legal': LEGAL
}

# Set poppler path - update this path according to your installation
POPPLER_PATH = r'C:\Program Files\poppler\Library\bin'
if os.path.exists(POPPLER_PATH):
    os.environ['PATH'] = POPPLER_PATH + os.pathsep + os.environ['PATH']
    logger.info(f'Poppler path added to PATH: {POPPLER_PATH}')
else:
    logger.warning('Poppler path not found. PDF to image conversion may not work.')

# Database initialization
def init_db():
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    c.execute('''
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            first_name TEXT NOT NULL,
            last_name TEXT NOT NULL,
            email TEXT UNIQUE NOT NULL,
            password TEXT NOT NULL,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    conn.commit()
    conn.close()

# Initialize database
init_db()

# JWT token required decorator
def token_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        token = None
        
        # Get token from header
        if 'Authorization' in request.headers:
            token = request.headers['Authorization'].split(" ")[1]
        
        if not token:
            return jsonify({'message': 'Token is missing'}), 401
        
        try:
            # Verify token
            data = jwt.decode(token, JWT_SECRET_KEY, algorithms=["HS256"])
            current_user = get_user_by_id(data['user_id'])
        except:
            return jsonify({'message': 'Token is invalid'}), 401
        
        return f(current_user, *args, **kwargs)
    return decorated

# User management functions
def get_user_by_id(user_id):
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    c.execute('SELECT id, first_name, last_name, email FROM users WHERE id = ?', (user_id,))
    user = c.fetchone()
    conn.close()
    
    if user:
        return {
            'id': user[0],
            'first_name': user[1],
            'last_name': user[2],
            'email': user[3]
        }
    return None

def get_user_by_email(email):
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    c.execute('SELECT * FROM users WHERE email = ?', (email,))
    user = c.fetchone()
    conn.close()
    return user

# Authentication routes
@app.route('/api/register', methods=['POST'])
def register():
    try:
        data = request.get_json()
        
        # Validate required fields
        required_fields = ['first_name', 'last_name', 'email', 'password']
        for field in required_fields:
            if not data.get(field):
                return jsonify({'error': f'{field} is required'}), 400
        
        # Check if user already exists
        conn = sqlite3.connect('users.db')
        c = conn.cursor()
        c.execute('SELECT * FROM users WHERE email = ?', (data['email'],))
        if c.fetchone():
            conn.close()
            return jsonify({'error': 'Email already registered'}), 400
        
        # Hash password
        hashed_password = generate_password_hash(data['password'])
        
        # Insert new user
        c.execute('''
            INSERT INTO users (first_name, last_name, email, password)
            VALUES (?, ?, ?, ?)
        ''', (data['first_name'], data['last_name'], data['email'], hashed_password))
        conn.commit()
        conn.close()
        
        return jsonify({'message': 'Registration successful'}), 201
        
    except Exception as e:
        logger.error(f'Registration error: {e}')
        return jsonify({'error': 'Registration failed'}), 500

@app.route('/api/login', methods=['POST'])
def login():
    data = request.get_json()
    
    # Validate required fields
    if not data.get('email') or not data.get('password'):
        return jsonify({'message': 'Email and password are required'}), 400
    
    # Get user
    user = get_user_by_email(data['email'])
    if not user:
        return jsonify({'message': 'Invalid email or password'}), 401
    
    # Verify password
    if not check_password_hash(user[4], data['password']):
        return jsonify({'message': 'Invalid email or password'}), 401
    
    # Generate token
    token = jwt.encode({
        'user_id': user[0],
        'exp': datetime.datetime.utcnow() + datetime.timedelta(days=30 if data.get('rememberMe') else 1)
    }, JWT_SECRET_KEY, algorithm="HS256")
    
    return jsonify({
        'token': token,
        'user': {
            'id': user[0],
            'first_name': user[1],
            'last_name': user[2],
            'email': user[3]
        }
    })

@app.route('/api/logout', methods=['POST'])
@token_required
def logout(current_user):
    # In a stateless JWT system, the client should remove the token
    return jsonify({'message': 'Logged out successfully'})

@app.route('/api/user/profile', methods=['GET'])
@token_required
def get_profile(current_user):
    return jsonify(current_user)

@app.route('/api/user/profile', methods=['PUT'])
@token_required
def update_profile(current_user):
    data = request.get_json()
    
    # Validate email if provided
    if data.get('email') and data['email'] != current_user['email']:
        if get_user_by_email(data['email']):
            return jsonify({'message': 'Email already registered'}), 400
    
    # Update user profile
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    try:
        c.execute('''
            UPDATE users 
            SET first_name = ?, last_name = ?, email = ?
            WHERE id = ?
        ''', (
            data.get('first_name', current_user['first_name']),
            data.get('last_name', current_user['last_name']),
            data.get('email', current_user['email']),
            current_user['id']
        ))
        conn.commit()
        return jsonify({'message': 'Profile updated successfully'})
    except Exception as e:
        return jsonify({'message': 'Profile update failed'}), 500
    finally:
        conn.close()

@app.route('/api/user/password', methods=['PUT'])
@token_required
def change_password(current_user):
    data = request.get_json()
    
    if not data.get('current_password') or not data.get('new_password'):
        return jsonify({'message': 'Current and new passwords are required'}), 400
    
    # Verify current password
    user = get_user_by_id(current_user['id'])
    if not check_password_hash(user[4], data['current_password']):
        return jsonify({'message': 'Current password is incorrect'}), 401
    
    # Update password
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    try:
        c.execute('''
            UPDATE users 
            SET password = ?
            WHERE id = ?
        ''', (generate_password_hash(data['new_password']), current_user['id']))
        conn.commit()
        return jsonify({'message': 'Password updated successfully'})
    except Exception as e:
        return jsonify({'message': 'Password update failed'}), 500
    finally:
        conn.close()

@app.after_request
def after_request(response):
    response.headers.add('Access-Control-Allow-Origin', '*')
    response.headers.add('Access-Control-Allow-Headers', 'Content-Type')
    response.headers.add('Access-Control-Allow-Methods', 'GET, POST, OPTIONS')
    return response

@app.route('/')
def index():
    try:
        logger.info('Serving index.html')
        return send_from_directory(BASE_DIR, 'index.html')
    except Exception as e:
        logger.error(f'Error serving index.html: {e}')
        abort(500)

@app.route('/register.html')
def register_page():
    try:
        logger.info('Serving register.html')
        return send_from_directory(BASE_DIR, 'register.html')
    except Exception as e:
        logger.error(f'Error serving register.html: {e}')
        abort(500)

@app.route('/convert-jpg-to-pdf', methods=['POST'])
def convert_jpg_to_pdf():
    try:
        logger.info('Starting JPG to PDF conversion')
        
        if 'files' not in request.files:
            logger.error('No file part in the request')
            return jsonify({'error': 'No files provided'}), 400
        
        files = request.files.getlist('files')
        if not files or files[0].filename == '':
            logger.error('No selected files')
            return jsonify({'error': 'No files selected'}), 400
        
        # Get settings
        page_size = request.form.get('page_size', 'a4')
        orientation = request.form.get('orientation', 'auto')
        
        logger.info(f'Processing {len(files)} images with page size: {page_size}, orientation: {orientation}')

        # Create temporary PDF file
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as pdf_temp:
            logger.info(f'Created temporary PDF file: {pdf_temp.name}')
            
            try:
                # Create PDF
                c = canvas.Canvas(pdf_temp.name, pagesize=PAGE_SIZES.get(page_size, A4))
                
                # Process each image
                for file in files:
                    if not file.filename.lower().endswith(('.jpg', '.jpeg')):
                        continue
                    
                    # Open and process image
                    img = Image.open(file.stream)
                    
                    # Determine page size and orientation
                    page_width, page_height = PAGE_SIZES.get(page_size, A4)
                    
                    if page_size == 'fit':
                        # Use image dimensions
                        page_width, page_height = img.size
                        c.setPageSize((page_width, page_height))
                    elif orientation == 'auto':
                        # Auto-detect orientation based on image dimensions
                        if (img.width > img.height) != (page_width > page_height):
                            page_width, page_height = page_height, page_width
                            c.setPageSize((page_width, page_height))
                    elif orientation == 'landscape' and page_width < page_height:
                        # Force landscape
                        page_width, page_height = page_height, page_width
                        c.setPageSize((page_width, page_height))
                    elif orientation == 'portrait' and page_width > page_height:
                        # Force portrait
                        page_width, page_height = page_height, page_width
                        c.setPageSize((page_width, page_height))
                    
                    # Calculate scaling to fit page while maintaining aspect ratio
                    scale_width = page_width / img.width
                    scale_height = page_height / img.height
                    scale = min(scale_width, scale_height)
                    
                    # Calculate centered position
                    width = img.width * scale
                    height = img.height * scale
                    x = (page_width - width) / 2
                    y = (page_height - height) / 2
                    
                    # Save image to temporary file
                    img_temp = tempfile.NamedTemporaryFile(suffix='.jpg', delete=False)
                    img.save(img_temp.name, 'JPEG')
                    
                    # Draw image on PDF
                    c.drawImage(img_temp.name, x, y, width=width, height=height)
                    c.showPage()
                    
                    # Clean up temporary image file
                    os.unlink(img_temp.name)
                    logger.info(f'Processed image: {file.filename}')
                
                c.save()
                logger.info('PDF creation completed successfully')

                # Send the PDF file
                return send_file(
                    pdf_temp.name,
                    as_attachment=True,
                    download_name='converted.pdf',
                    mimetype='application/pdf'
                )

            except Exception as conv_error:
                logger.error(f'Error during conversion: {str(conv_error)}')
                raise
            finally:
                # Clean up temporary file
                logger.info('Cleaning up temporary files')
                try:
                    os.unlink(pdf_temp.name)
                except Exception as cleanup_error:
                    logger.error(f'Error cleaning up temporary files: {str(cleanup_error)}')

    except Exception as e:
        error_msg = str(e)
        logger.error(f'Error converting JPG to PDF: {error_msg}')
        return jsonify({
            'error': 'Failed to convert JPG to PDF',
            'details': error_msg
        }), 500

@app.route('/convert-pdf-to-jpeg', methods=['POST'])
def convert_pdf_to_jpeg():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '' or not file.filename.lower().endswith('.pdf'):
            return jsonify({'error': 'Invalid file format. Please upload a PDF file.'}), 400

        # Get conversion settings
        quality = int(request.form.get('quality', 300))
        start_page = int(request.form.get('start_page', 1))
        end_page = request.form.get('end_page')
        if end_page:
            end_page = int(end_page)

        # Create a temporary directory for processing
        with tempfile.TemporaryDirectory() as temp_dir:
            pdf_path = os.path.join(temp_dir, secure_filename(file.filename))
            file.save(pdf_path)

            # Convert PDF to images using pdf2image
            try:
                images = convert_from_path(
                    pdf_path,
                    dpi=quality,
                    first_page=start_page,
                    last_page=end_page
                )
            except Exception as e:
                return jsonify({'error': f'PDF conversion failed: {str(e)}'}), 500

            # Create a ZIP file containing all converted images
            zip_path = os.path.join(temp_dir, 'converted_images.zip')
            with zipfile.ZipFile(zip_path, 'w') as zip_file:
                for i, image in enumerate(images, start=start_page):
                    image_path = os.path.join(temp_dir, f'page_{i}.jpg')
                    image.save(image_path, 'JPEG', quality=95)
                    zip_file.write(image_path, f'page_{i}.jpg')

            return send_file(
                zip_path,
                mimetype='application/zip',
                as_attachment=True,
                download_name=f"{os.path.splitext(file.filename)[0]}_images.zip"
            )

    except Exception as e:
        app.logger.error(f"Error in PDF to JPEG conversion: {str(e)}")
        return jsonify({'error': 'An error occurred during conversion'}), 500

@app.route('/convert-pdf-to-word', methods=['POST'])
def convert_pdf_to_word():
    try:
        logger.info('Starting PDF to Word conversion')
        
        if 'file' not in request.files:
            logger.error('No file part in the request')
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '':
            logger.error('No selected file')
            return jsonify({'error': 'No file selected'}), 400
        
        if not file.filename.lower().endswith('.pdf'):
            logger.error(f'Invalid file type: {file.filename}')
            return jsonify({'error': 'Only PDF files are allowed'}), 400

        logger.info(f'Processing file: {file.filename}')

        # Create temporary files for processing
        try:
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as pdf_temp:
                logger.info(f'Created temporary PDF file: {pdf_temp.name}')
                file.save(pdf_temp.name)
                
                docx_temp = tempfile.NamedTemporaryFile(suffix='.docx', delete=False)
                docx_temp.close()
                logger.info(f'Created temporary DOCX file: {docx_temp.name}')

                try:
                    # Convert PDF to DOCX
                    logger.info('Starting PDF to DOCX conversion')
                    cv = Converter(pdf_temp.name)
                    cv.convert(docx_temp.name)
                    cv.close()
                    logger.info('PDF to DOCX conversion completed successfully')

                    # Verify the output file exists and has content
                    if not os.path.exists(docx_temp.name):
                        raise FileNotFoundError('Converted DOCX file not found')
                    
                    if os.path.getsize(docx_temp.name) == 0:
                        raise ValueError('Converted DOCX file is empty')

                    # Send the converted file
                    logger.info('Sending converted file to client')
                    return send_file(
                        docx_temp.name,
                        as_attachment=True,
                        download_name='converted.docx',
                        mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                    )
                except Exception as conv_error:
                    logger.error(f'Error during conversion: {str(conv_error)}')
                    raise
                finally:
                    # Clean up temporary files
                    logger.info('Cleaning up temporary files')
                    try:
                        os.unlink(pdf_temp.name)
                        os.unlink(docx_temp.name)
                    except Exception as cleanup_error:
                        logger.error(f'Error cleaning up temporary files: {str(cleanup_error)}')

        except Exception as temp_error:
            logger.error(f'Error handling temporary files: {str(temp_error)}')
            raise

    except Exception as e:
        error_msg = str(e)
        logger.error(f'Error converting PDF to Word: {error_msg}')
        return jsonify({
            'error': 'Failed to convert PDF to Word',
            'details': error_msg
        }), 500

@app.route('/convert-word-to-pdf', methods=['POST'])
def convert_word_to_pdf():
    try:
        if 'file' not in request.files:
            return 'No file uploaded', 400

        file = request.files['file']
        if file.filename == '':
            return 'No file selected', 400

        if not file.filename.lower().endswith(('.doc', '.docx')):
            return 'Invalid file format. Please upload a Word document (.doc or .docx)', 400

        # Create temporary files for processing
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_word:
            file.save(temp_word.name)
            
            # Create temporary file for PDF output
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
                try:
                    # Load the Word document
                    doc = Document(temp_word.name)
                    
                    # Create PDF
                    c = canvas.Canvas(temp_pdf.name, pagesize=LETTER)
                    width, height = LETTER
                    y = height - inch  # Start from top of page
                    
                    # Process each paragraph
                    for para in doc.paragraphs:
                        if y < inch:  # Check if we need a new page
                            c.showPage()
                            y = height - inch
                        
                        # Add text to PDF
                        c.drawString(inch, y, para.text)
                        y -= 14  # Move down for next line
                    
                    c.save()
                    
                    # Return the PDF file
                    return send_file(
                        temp_pdf.name,
                        as_attachment=True,
                        download_name=secure_filename(file.filename.rsplit('.', 1)[0] + '.pdf'),
                        mimetype='application/pdf'
                    )
                
                except Exception as e:
                    logging.error(f"Error converting Word to PDF: {str(e)}")
                    return f"Error converting file: {str(e)}", 500
                
                finally:
                    # Clean up temporary files
                    try:
                        os.unlink(temp_word.name)
                        os.unlink(temp_pdf.name)
                    except Exception as e:
                        logging.error(f"Error cleaning up temporary files: {str(e)}")

    except Exception as e:
        logging.error(f"Error processing Word to PDF conversion: {str(e)}")
        return f"Error processing file: {str(e)}", 500

@app.route('/convert-powerpoint-to-pdf', methods=['POST'])
def convert_powerpoint_to_pdf():
    try:
        if 'file' not in request.files:
            return 'No file uploaded', 400

        file = request.files['file']
        if file.filename == '':
            return 'No file selected', 400

        if not file.filename.lower().endswith(('.ppt', '.pptx')):
            return 'Invalid file format. Please upload a PowerPoint presentation (.ppt or .pptx)', 400

        # Create temporary files for processing
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_ppt:
            file.save(temp_ppt.name)
            
            # Create temporary file for PDF output
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
                try:
                    # Load the PowerPoint presentation
                    prs = Presentation(temp_ppt.name)
                    
                    # Create PDF
                    c = canvas.Canvas(temp_pdf.name, pagesize=LETTER)
                    width, height = letter
                    
                    # Process each slide
                    for slide in prs.slides:
                        # Start a new page for each slide
                        y = height - inch
                        
                        # Process shapes in the slide (text boxes, etc.)
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                # Draw text from the shape
                                text = shape.text.strip()
                                if text:
                                    c.drawString(inch, y, text)
                                    y -= 14  # Move down for next line
                        
                        c.showPage()  # Move to next page
                    
                    c.save()
                    
                    # Return the PDF file
                    return send_file(
                        temp_pdf.name,
                        as_attachment=True,
                        download_name=secure_filename(file.filename.rsplit('.', 1)[0] + '.pdf'),
                        mimetype='application/pdf'
                    )
                
                except Exception as e:
                    logging.error(f"Error converting PowerPoint to PDF: {str(e)}")
                    return f"Error converting file: {str(e)}", 500
                
                finally:
                    # Clean up temporary files
                    try:
                        os.unlink(temp_ppt.name)
                        os.unlink(temp_pdf.name)
                    except Exception as e:
                        logging.error(f"Error cleaning up temporary files: {str(e)}")

    except Exception as e:
        logging.error(f"Error processing PowerPoint to PDF conversion: {str(e)}")
        return f"Error processing file: {str(e)}", 500

@app.route('/convert-excel-to-pdf', methods=['POST'])
def convert_excel_to_pdf():
    try:
        if 'file' not in request.files:
            return 'No file uploaded', 400

        file = request.files['file']
        if file.filename == '':
            return 'No file selected', 400

        if not file.filename.lower().endswith(('.xls', '.xlsx')):
            return 'Invalid file format. Please upload an Excel file (.xls or .xlsx)', 400

        # Get conversion options
        fit_to_page = request.form.get('fitToPage', 'true').lower() == 'true'
        include_gridlines = request.form.get('includeGridlines', 'true').lower() == 'true'
        all_worksheets = request.form.get('allWorksheets', 'true').lower() == 'true'

        # Create temporary files for processing
        with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as temp_excel:
            file.save(temp_excel.name)
            
            # Create temporary file for PDF output
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
                try:
                    # Load the Excel workbook
                    wb = load_workbook(temp_excel.name)
                    
                    # Create PDF
                    c = canvas.Canvas(temp_pdf.name, pagesize=letter)
                    width, height = letter
                    
                    # Process each worksheet
                    for ws in wb.worksheets:
                        if not all_worksheets and ws != wb.active:
                            continue
                            
                        # Add worksheet name as header
                        c.setFont("Helvetica-Bold", 14)
                        c.drawString(inch, height - inch, f"Sheet: {ws.title}")
                        c.setFont("Helvetica", 10)
                        
                        # Calculate column widths
                        col_widths = []
                        for col in ws.columns:
                            max_length = 0
                            for cell in col:
                                try:
                                    if cell.value:
                                        max_length = max(max_length, len(str(cell.value)))
                                except:
                                    pass
                            col_widths.append(min(max_length * 7, 100))  # Scale factor and max width
                        
                        # Draw table
                        y = height - 2*inch
                        x = inch
                        row_height = 20
                        
                        # Draw headers and data
                        for row_idx, row in enumerate(ws.rows):
                            if y < inch:  # Start new page if needed
                                c.showPage()
                                y = height - inch
                                c.setFont("Helvetica-Bold", 14)
                                c.drawString(inch, y, f"Sheet: {ws.title} (continued)")
                                c.setFont("Helvetica", 10)
                                y -= row_height
                            
                            x = inch
                            for col_idx, cell in enumerate(row):
                                # Draw cell border if gridlines are enabled
                                if include_gridlines:
                                    c.rect(x, y, col_widths[col_idx], row_height)
                                
                                # Draw cell content
                                if cell.value is not None:
                                    c.drawString(x + 2, y + 5, str(cell.value))
                                
                                x += col_widths[col_idx]
                            
                            y -= row_height
                        
                        # Start new page for next worksheet
                        c.showPage()
                    
                    c.save()
                    
                    # Return the PDF file
                    return send_file(
                        temp_pdf.name,
                        as_attachment=True,
                        download_name=secure_filename(file.filename.rsplit('.', 1)[0] + '.pdf'),
                        mimetype='application/pdf'
                    )
                
                except Exception as e:
                    logging.error(f"Error converting Excel to PDF: {str(e)}")
                    return f"Error converting file: {str(e)}", 500
                
                finally:
                    # Clean up temporary files
                    try:
                        os.unlink(temp_excel.name)
                        os.unlink(temp_pdf.name)
                    except Exception as e:
                        logging.error(f"Error cleaning up temporary files: {str(e)}")

    except Exception as e:
        logging.error(f"Error processing Excel to PDF conversion: {str(e)}")
        return f"Error processing file: {str(e)}", 500

@app.route('/convert-html-to-pdf', methods=['POST'])
def convert_html_to_pdf():
    try:
        # Get conversion options
        include_styles = request.form.get('includeStyles', 'true').lower() == 'true'
        include_images = request.form.get('includeImages', 'true').lower() == 'true'
        page_breaks = request.form.get('pageBreaks', 'true').lower() == 'true'

        # Configure pdfkit options
        options = {
            'enable-local-file-access': None,
            'load-error-handling': 'ignore',
            'load-media-error-handling': 'ignore'
        }

        if not include_styles:
            options['no-default-style'] = None
        if not include_images:
            options['no-images'] = None
        if page_breaks:
            options['page-size'] = 'Letter'

        # Create temporary file for PDF output
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            try:
                if 'file' in request.files:
                    # Handle file upload
                    file = request.files['file']
                    if file.filename == '':
                        return 'No file selected', 400

                    if not file.filename.lower().endswith(('.html', '.htm')):
                        return 'Invalid file format. Please upload an HTML file (.html or .htm)', 400

                    # Save uploaded file temporarily
                    with tempfile.NamedTemporaryFile(delete=False, suffix='.html') as temp_html:
                        file.save(temp_html.name)
                        pdfkit.from_file(temp_html.name, temp_pdf.name, options=options)
                        os.unlink(temp_html.name)

                elif 'url' in request.form:
                    # Handle URL input
                    url = request.form['url']
                    try:
                        # Validate URL by making a HEAD request
                        response = requests.head(url, allow_redirects=True)
                        response.raise_for_status()
                        pdfkit.from_url(url, temp_pdf.name, options=options)
                    except requests.exceptions.RequestException as e:
                        return f'Error accessing URL: {str(e)}', 400

                elif 'html' in request.form:
                    # Handle direct HTML input
                    html_content = request.form['html']
                    if not html_content.strip():
                        return 'No HTML content provided', 400
                    pdfkit.from_string(html_content, temp_pdf.name, options=options)

                else:
                    return 'No input provided', 400

                # Return the PDF file
                return send_file(
                    temp_pdf.name,
                    as_attachment=True,
                    download_name='converted.pdf',
                    mimetype='application/pdf'
                )

            except Exception as e:
                logging.error(f"Error converting HTML to PDF: {str(e)}")
                return f"Error converting file: {str(e)}", 500

            finally:
                # Clean up temporary file
                try:
                    os.unlink(temp_pdf.name)
                except Exception as e:
                    logging.error(f"Error cleaning up temporary files: {str(e)}")

    except Exception as e:
        logging.error(f"Error processing HTML to PDF conversion: {str(e)}")
        return f"Error processing file: {str(e)}", 500

@app.route('/<path:path>')
def serve_file(path):
    try:
        logger.info(f'Attempting to serve file: {path}')
        
        # Map common directories to their locations
        if path.startswith(('components/', 'js/', 'css/', 'tools/')):
            directory = os.path.join(BASE_DIR, os.path.dirname(path))
            filename = os.path.basename(path)
            full_path = os.path.join(directory, filename)
            
            if os.path.exists(full_path):
                logger.info(f'Serving file from {directory}: {filename}')
                return send_from_directory(directory, filename)
        
        # Try serving from root directory
        full_path = os.path.join(BASE_DIR, path)
        if os.path.exists(full_path):
            logger.info(f'Serving file from root: {path}')
            return send_from_directory(BASE_DIR, path)
            
        logger.error(f'File not found: {path} (full path: {full_path})')
        abort(404)
    except Exception as e:
        logger.error(f'Error serving {path}: {e}')
        abort(404)

@app.errorhandler(404)
def not_found(e):
    logger.error(f'404 error: {e}')
    return f"File not found: {e}", 404

@app.errorhandler(500)
def server_error(e):
    logger.error(f'500 error: {e}')
    return f"Internal server error: {e}", 500

@app.route('/add-watermark', methods=['POST'])
def add_watermark():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400
        
        file = request.files['file']
        if not file or not file.filename.lower().endswith('.pdf'):
            return jsonify({'error': 'Invalid file format. Please upload a PDF file.'}), 400

        # Create temporary files
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as pdf_temp:
            file.save(pdf_temp.name)
            pdf_path = pdf_temp.name

        watermark_type = request.form.get('watermark_type', 'text')
        position = request.form.get('position', 'center')
        rotation = int(request.form.get('rotation', 0))
        page_range = request.form.get('page_range', '')

        # Create output PDF
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as output_temp:
            output_path = output_temp.name

        # Read the PDF
        pdf_reader = PyPDF2.PdfReader(pdf_path)
        pdf_writer = PyPDF2.PdfWriter()

        # Process each page
        for page_num in range(len(pdf_reader.pages)):
            # Check if this page should be watermarked
            if page_range:
                page_ranges = parse_page_range(page_range)
                if page_num + 1 not in page_ranges:
                    pdf_writer.add_page(pdf_reader.pages[page_num])
                    continue

            page = pdf_reader.pages[page_num]
            
            # Create a new page with watermark
            packet = io.BytesIO()
            can = canvas.Canvas(packet, pagesize=page.mediabox)
            
            # Get page dimensions
            width = float(page.mediabox.width)
            height = float(page.mediabox.height)

            if watermark_type == 'text':
                # Get text watermark parameters
                text = request.form.get('watermark_text', '')
                font_size = int(request.form.get('font_size', 16))
                font_color = request.form.get('font_color', '#000000')
                opacity = float(request.form.get('opacity', 50)) / 100

                # Set font and color
                can.setFont("Helvetica", font_size)
                r, g, b = hex_to_rgb(font_color)
                can.setFillColorRGB(r, g, b, alpha=opacity)

                # Calculate text position based on position parameter
                text_width = can.stringWidth(text, "Helvetica", font_size)
                text_height = font_size

                if position == 'center':
                    x = (width - text_width) / 2
                    y = (height - text_height) / 2
                elif position == 'top-left':
                    x = 50
                    y = height - 50
                elif position == 'top-right':
                    x = width - text_width - 50
                    y = height - 50
                elif position == 'bottom-left':
                    x = 50
                    y = 50
                else:  # bottom-right
                    x = width - text_width - 50
                    y = 50

                # Save the current state
                can.saveState()
                # Translate to center and rotate
                can.translate(x + text_width/2, y + text_height/2)
                can.rotate(rotation)
                can.translate(-(x + text_width/2), -(y + text_height/2))
                # Draw text
                can.drawString(x, y, text)
                # Restore the state
                can.restoreState()

            else:  # image watermark
                if 'watermark_image' not in request.files:
                    return jsonify({'error': 'No watermark image uploaded'}), 400

                image_file = request.files['watermark_image']
                image_opacity = float(request.form.get('image_opacity', 50)) / 100

                # Save watermark image temporarily
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as img_temp:
                    image_file.save(img_temp.name)
                    img_path = img_temp.name

                # Get image dimensions
                img = Image.open(img_path)
                img_width, img_height = img.size

                # Calculate image position and size
                max_size = min(width, height) * 0.5  # 50% of page size
                scale = min(max_size / img_width, max_size / img_height)
                new_width = img_width * scale
                new_height = img_height * scale

                if position == 'center':
                    x = (width - new_width) / 2
                    y = (height - new_height) / 2
                elif position == 'top-left':
                    x = 50
                    y = height - new_height - 50
                elif position == 'top-right':
                    x = width - new_width - 50
                    y = height - new_height - 50
                elif position == 'bottom-left':
                    x = 50
                    y = 50
                else:  # bottom-right
                    x = width - new_width - 50
                    y = 50

                # Draw image with rotation
                can.saveState()
                can.translate(x + new_width/2, y + new_height/2)
                can.rotate(rotation)
                can.translate(-(x + new_width/2), -(y + new_height/2))
                can.drawImage(img_path, x, y, width=new_width, height=new_height)
                can.restoreState()

                # Clean up temporary image file
                os.unlink(img_path)

            can.save()
            packet.seek(0)

            # Create a new PDF with the watermark
            new_pdf = PyPDF2.PdfReader(packet)
            page.merge_page(new_pdf.pages[0])
            pdf_writer.add_page(page)

        # Save the watermarked PDF
        with open(output_path, 'wb') as output_file:
            pdf_writer.write(output_file)

        # Clean up temporary files
        os.unlink(pdf_path)

        # Return the watermarked PDF
        return send_file(
            output_path,
            mimetype='application/pdf',
            as_attachment=True,
            download_name='watermarked.pdf'
        )

    except Exception as e:
        logging.error(f"Error adding watermark: {str(e)}")
        return jsonify({'error': 'An error occurred while adding watermark'}), 500

    finally:
        # Clean up output file
        if 'output_path' in locals():
            try:
                os.unlink(output_path)
            except:
                pass

def hex_to_rgb(hex_color):
    """Convert hex color to RGB values."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) / 255 for i in (0, 2, 4))

def parse_page_range(page_range):
    """Parse page range string into list of page numbers."""
    if not page_range:
        return []
    
    pages = set()
    parts = page_range.split(',')
    
    for part in parts:
        if '-' in part:
            start, end = map(int, part.split('-'))
            pages.update(range(start, end + 1))
        else:
            pages.add(int(part))
    
    return sorted(list(pages))

@app.route('/rotate-pdf', methods=['POST'])
def rotate_pdf():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400
        
        file = request.files['file']
        if not file or not file.filename.lower().endswith('.pdf'):
            return jsonify({'error': 'Invalid file format. Please upload a PDF file.'}), 400

        # Get rotation parameters
        rotation_angle = int(request.form.get('rotation_angle', 90))
        rotation_scope = request.form.get('rotation_scope', 'all')
        page_range = request.form.get('page_range', '')

        # Create temporary files
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as pdf_temp:
            file.save(pdf_temp.name)
            pdf_path = pdf_temp.name

        # Create output PDF
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as output_temp:
            output_path = output_temp.name

        # Read the PDF
        pdf_reader = PyPDF2.PdfReader(pdf_path)
        pdf_writer = PyPDF2.PdfWriter()

        # Process each page
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            
            # Check if this page should be rotated
            if rotation_scope == 'specific' and page_range:
                page_ranges = parse_page_range(page_range)
                if page_num + 1 in page_ranges:
                    # Rotate the page
                    page.rotate(rotation_angle)
            
            pdf_writer.add_page(page)

        # Save the rotated PDF
        with open(output_path, 'wb') as output_file:
            pdf_writer.write(output_file)

        # Clean up temporary files
        os.unlink(pdf_path)

        # Return the rotated PDF
        return send_file(
            output_path,
            mimetype='application/pdf',
            as_attachment=True,
            download_name='rotated.pdf'
        )

    except Exception as e:
        logging.error(f"Error rotating PDF: {str(e)}")
        return jsonify({'error': 'An error occurred while rotating PDF'}), 500

    finally:
        # Clean up output file
        if 'output_path' in locals():
            try:
                os.unlink(output_path)
            except:
                pass

@app.route('/add-page-numbers', methods=['POST'])
def add_page_numbers():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400
        
        file = request.files['file']
        if not file or not file.filename.lower().endswith('.pdf'):
            return jsonify({'error': 'Invalid file format. Please upload a PDF file.'}), 400

        # Get page number parameters
        number_style = request.form.get('number_style', '1')
        number_position = request.form.get('number_position', 'bottom-right')
        page_range = request.form.get('page_range', '')
        font_size = int(request.form.get('font_size', 10))
        font_color = request.form.get('font_color', '#000000')

        # Create temporary files
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as pdf_temp:
            file.save(pdf_temp.name)
            pdf_path = pdf_temp.name

        # Create output PDF
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as output_temp:
            output_path = output_temp.name

        # Read the PDF
        pdf_reader = PyPDF2.PdfReader(pdf_path)
        pdf_writer = PyPDF2.PdfWriter()

        # Process each page
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            
            # Check if this page should be numbered
            if page_range:
                page_ranges = parse_page_range(page_range)
                if page_num + 1 not in page_ranges:
                    pdf_writer.add_page(page)
                    continue

            # Create a new page with page number
            packet = io.BytesIO()
            can = canvas.Canvas(packet, pagesize=page.mediabox)
            
            # Get page dimensions
            width = float(page.mediabox.width)
            height = float(page.mediabox.height)

            # Set font and color
            can.setFont("Helvetica", font_size)
            r, g, b = hex_to_rgb(font_color)
            can.setFillColorRGB(r, g, b)

            # Convert page number to desired style
            page_number = page_num + 1
            if number_style == 'i':
                page_number = roman.toRoman(page_number).lower()
            elif number_style == 'I':
                page_number = roman.toRoman(page_number)
            elif number_style == 'a':
                page_number = chr(96 + page_number)  # a=97, b=98, etc.
            elif number_style == 'A':
                page_number = chr(64 + page_number)  # A=65, B=66, etc.

            # Calculate text position based on position parameter
            text = str(page_number)
            text_width = can.stringWidth(text, "Helvetica", font_size)
            text_height = font_size

            # Set position based on parameter
            if number_position == 'bottom-right':
                x = width - text_width - 50
                y = 50
            elif number_position == 'bottom-center':
                x = (width - text_width) / 2
                y = 50
            elif number_position == 'bottom-left':
                x = 50
                y = 50
            elif number_position == 'top-right':
                x = width - text_width - 50
                y = height - 50
            elif number_position == 'top-center':
                x = (width - text_width) / 2
                y = height - 50
            else:  # top-left
                x = 50
                y = height - 50

            # Draw page number
            can.drawString(x, y, text)
            can.save()
            packet.seek(0)

            # Create a new PDF with the page number
            new_pdf = PyPDF2.PdfReader(packet)
            page.merge_page(new_pdf.pages[0])
            pdf_writer.add_page(page)

        # Save the numbered PDF
        with open(output_path, 'wb') as output_file:
            pdf_writer.write(output_file)

        # Clean up temporary files
        os.unlink(pdf_path)

        # Return the numbered PDF
        return send_file(
            output_path,
            mimetype='application/pdf',
            as_attachment=True,
            download_name='numbered.pdf'
        )

    except Exception as e:
        logging.error(f"Error adding page numbers: {str(e)}")
        return jsonify({'error': 'An error occurred while adding page numbers'}), 500

    finally:
        # Clean up output file
        if 'output_path' in locals():
            try:
                os.unlink(output_path)
            except:
                pass

if __name__ == '__main__':
    try:
        port = 8080
        logger.info(f'Starting Flask server on port {port}...')
        logger.info(f'Base directory: {BASE_DIR}')
        app.run(debug=True, host='0.0.0.0', port=port)
    except Exception as e:
        logger.error(f'Error starting server: {e}')
        raise 